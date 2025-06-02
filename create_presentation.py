from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(59, 89, 152)  # Dark blue
    SECONDARY_COLOR = RGBColor(66, 133, 244)  # Light blue
    ACCENT_COLOR = RGBColor(255, 153, 0)  # Orange
    WHITE = RGBColor(255, 255, 255)
    
    # ===== SLIDE 1: TITLE =====
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title = slide.shapes.title
    title.text = "Office Archive Management System"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Solusi Digital untuk Manajemen Dokumen Perkantoran yang Efisien"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(1))
    tf = footer.text_frame
    p = tf.add_paragraph()
    p.text = "Presentasi - " + datetime.now().strftime("%d %B %Y")
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # ===== SLIDE 2: AGENDA =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Agenda"
    
    content = [
        "Gambaran Umum Sistem",
        "Fitur Utama Aplikasi",
        "Teknologi yang Digunakan",
        "Demo Aplikasi",
        "Keunggulan Solusi",
        "Rencana Pengembangan"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for item in content:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 3: GAMBARAN UMUM =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Gambaran Umum"
    
    points = [
        "Platform manajemen dokumen digital terintegrasi",
        "Menyediakan solusi pengarsipan yang terstruktur",
        "Mendukung berbagai format dokumen",
        "Memudahkan pelacakan dan pencarian dokumen",
        "Mengotomatiskan alur kerja persetujuan dokumen"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 4: FITUR UTAMA 1 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fitur Utama - Manajemen Dokumen"
    
    points = [
        "Upload berbagai format file (PDF, DOC, XLS, dll)",
        "Pencarian canggih dengan filter",
        "Preview dokumen langsung di browser",
        "Riwayat perubahan dan versi dokumen",
        "Manajemen kategori terstruktur"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 5: FITUR UTAMA 2 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fitur Utama - SPK & PPBJ"
    
    points = [
        "Form pengajuan online yang mudah digunakan",
        "Proses persetujuan multi-level",
        "Notifikasi real-time",
        "Pelacakan status pengajuan",
        "Arsip digital dokumen yang disetujui"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 6: TEKNOLOGI =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Teknologi yang Digunakan"
    
    tech = [
        ("Backend", "PHP 8.1+, Laravel 10.x, MySQL 8.0+"),
        ("Frontend", "Bootstrap 5, jQuery, DataTables"),
        ("Keamanan", "Autentikasi multi-faktor, Enkripsi data, Audit log"),
        ("Hosting", "Dapat di-host di berbagai platform cloud")
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for category, items in tech:
        p = tf.add_paragraph()
        p.text = f"{category}:"
        p.level = 0
        p.font.bold = True
        p.font.size = Pt(22)
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = items
        p.level = 1
        p.font.size = Pt(20)
        p.space_after = Pt(16)
    
    # ===== SLIDE 7: DEMO =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Demo Aplikasi"
    
    points = [
        "Login dan Dashboard",
        "Manajemen Dokumen",
        "Pengajuan SPK & PPBJ",
        "Manajemen Pengguna",
        "Laporan dan Analitik"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = Pt(28)
        p.space_after = Pt(12)
    
    # ===== SLIDE 8: KEUNGGULAN =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Keunggulan Solusi"
    
    points = [
        "Efisiensi waktu dengan pengelolaan dokumen terpusat",
        "Akses mudah kapan saja, di mana saja",
        "Keamanan data terjamin dengan enkripsi",
        "Antarmuka yang ramah pengguna",
        "Dapat disesuaikan dengan kebutuhan bisnis"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 9: RENCANA PENGEMBANGAN =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Rencana Pengembangan"
    
    points = [
        "Integrasi dengan layanan cloud storage",
        "Pengembangan aplikasi mobile",
        "Fitur tanda tangan digital",
        "Analitik dan pelaporan yang lebih canggih",
        "Integrasi dengan sistem eksternal"
    ]
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    # ===== SLIDE 10: PENUTUP =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Terima Kasih"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = "Ada pertanyaan?"
    p.font.size = Pt(32)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Add contact info
    contact_info = [
        ("📧", "Email", "info@perusahaananda.com"),
        ("📱", "Telepon/WA", "+62 822-8721-8274"),
        ("🌐", "Website", "www.perusahaananda.com")
    ]
    
    y_pos = 4.0
    for icon, label, value in contact_info:
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.add_paragraph()
        p.text = icon
        p.font.size = Pt(20)
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(4.7), Inches(y_pos), Inches(2), Inches(0.5))
        tf = label_box.text_frame
        p = tf.add_paragraph()
        p.text = f"{label}:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(6), Inches(y_pos), Inches(4), Inches(0.5))
        tf = value_box.text_frame
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(18)
        
        y_pos += 0.6
    
    # Add footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    tf = footer.text_frame
    p = tf.add_paragraph()
    p.text = "© 2024 Office Archive Management System. All rights reserved."
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    prs.save('presentasi_office_archive.pptx')
    print("Presentasi berhasil dibuat: presentasi_office_archive.pptx")

if __name__ == "__main__":
    create_presentation()
